from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Theme (matches the app: ppt dark theme)
BG = RGBColor(0x1E, 0x1E, 0x1E)
CARD = RGBColor(0x1D, 0x1D, 0x1D)
PRIMARY = RGBColor(0x00, 0xB8, 0xEC)   # cyan
ACCENT = RGBColor(0xEC, 0x6A, 0x37)    # orange
SUCCESS = RGBColor(0x9E, 0xCD, 0x3B)   # green
TEXT = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xBD, 0xC1, 0xC6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, 18, TEXT, False)]
    first = True
    for item in runs:
        text, size, color, bold = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = space
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
        r.font.name = "Arial"
    return tb


def bar(slide):
    box(slide, 0, 0, SW, Inches(0.18), fill=PRIMARY)
    box(slide, 0, Inches(0.18), Inches(4.4), Emu(int(Inches(0.06))), fill=ACCENT)


def header(slide, kicker, title):
    bar(slide)
    txt(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.4),
        [(kicker.upper(), 13, ACCENT, True)])
    txt(slide, Inches(0.6), Inches(0.8), Inches(12.1), Inches(0.9),
        [(title, 30, TEXT, True)])


def bullets(slide, x, y, w, h, items, size=16, gap=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        lead = it[0] if isinstance(it, tuple) else None
        text = it[1] if isinstance(it, tuple) else it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = gap
        r = p.add_run(); r.text = "▸  "
        r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = PRIMARY; r.font.name = "Arial"
        if lead:
            r2 = p.add_run(); r2.text = lead + "  "
            r2.font.size = Pt(size); r2.font.bold = True
            r2.font.color.rgb = TEXT; r2.font.name = "Arial"
        r3 = p.add_run(); r3.text = text
        r3.font.size = Pt(size); r3.font.color.rgb = MUTED; r3.font.name = "Arial"
    return tb


def card(slide, x, y, w, h, title, body, accent=PRIMARY):
    box(slide, x, y, w, h, fill=CARD, line=RGBColor(0x3A, 0x3A, 0x3A), line_w=1)
    box(slide, x, y, Emu(int(Inches(0.08))), h, fill=accent)
    txt(slide, x + Inches(0.25), y + Inches(0.18), w - Inches(0.4), Inches(0.5),
        [(title, 16, TEXT, True)])
    txt(slide, x + Inches(0.25), y + Inches(0.72), w - Inches(0.4), h - Inches(0.85),
        [(body, 12.5, MUTED, False)], space=1.1)


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0, 0, SW, Inches(0.22), fill=PRIMARY)
box(s, 0, Inches(0.22), Inches(5), Emu(int(Inches(0.08))), fill=ACCENT)
txt(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.2),
    [("Skintyee First Nation", 22, PRIMARY, True)])
txt(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.6),
    [("Community Mobile App", 52, TEXT, True)])
txt(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.7),
    [("Proposal & Proof-of-Concept", 22, ACCENT, True)])
txt(s, Inches(0.92), Inches(5.2), Inches(11.5), Inches(0.6),
    [("React Native · Azure · One app for Public, Members, and Band Staff", 15, MUTED, False)])
txt(s, Inches(0.92), Inches(6.7), Inches(11.5), Inches(0.4),
    [("Prepared by Blue Collar Dev  ·  3-month engagement  ·  May–August 2026", 12, MUTED, False)])

# ---------- Slide 2: Overview ----------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Overview", "One app, three audiences")
bullets(s, Inches(0.7), Inches(1.95), Inches(7.2), Inches(5), [
    ("Goal:", "a single mobile + web app for the Skintyee First Nation community."),
    ("Public:", "community events and public records, open to anyone."),
    ("Band Members:", "directory, band meetings, polling & voting on issues."),
    ("Admins / Staff:", "time keeping and financial records, admin-only."),
    ("Built on Azure:", "identity, storage, and database all in Microsoft Azure."),
    ("Familiar feel:", "reuses our proven React Native app stack and design."),
], size=17, gap=1.5)
card(s, Inches(8.2), Inches(2.0), Inches(4.5), Inches(4.6),
     "Why now",
     "The current presence is on Site123. This app gives the Nation a modern, "
     "self-owned platform — role-based access, real-time community engagement, "
     "and auditable records — fully under Azure with easy backups.\n\n"
     "As an NGO, the priority is reliability and auditability over cost.",
     accent=ACCENT)

# ---------- Slide 3: Features (from diagram) ----------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Features", "What the app does")
feats = [
    ("Dashboard", "Community stats + spending charts at a glance.", SUCCESS),
    ("Transparency", "Band expenditures by area; drill into how much & where.", SUCCESS),
    ("Band Member Directory", "Searchable directory; contact details for members.", PRIMARY),
    ("Community Events", "Public event listings with details and dates.", PRIMARY),
    ("Notifications", "Health, Safety, Council… synced to skintyee.ca categories.", ACCENT),
    ("Band Meetings", "Agendas, schedules, and minutes for members.", PRIMARY),
    ("Time Keeping", "Hours logging & approval for workers (staff).", ACCENT),
    ("Financial Records", "Budgets and statements (admin only).", ACCENT),
    ("Polling + Surveys", "Vote on issues; live results.", PRIMARY),
    ("Auto-Publish", "Meetings & events pushed to skintyee.ca.", SUCCESS),
]
x0, y0 = Inches(0.7), Inches(1.95)
cw, ch, gx, gy = Inches(2.95), Inches(1.62), Inches(0.18), Inches(0.2)
for i, (t, b, c) in enumerate(feats):
    col = i % 4; row = i // 4
    card(s, x0 + col * (cw + gx), y0 + row * (ch + gy), cw, ch, t, b, accent=c)
txt(s, Inches(0.7), Inches(5.05), Inches(12), Inches(0.4),
    [("Access is gated by role — Public / Member / Admin+Staff — per the architecture diagram.", 12, MUTED, False)])

# ---------- Slide 3b: Transparency & Dashboard ----------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Transparency", "Open books, by design")
bullets(s, Inches(0.7), Inches(2.0), Inches(7.2), Inches(4.6), [
    ("Where the money goes:", "public band expenditures by program area."),
    ("Areas:", "Housing, Public Works, Education, Employment, Health, Social"),
    ("", "Assistance, Child & Family Services, IT, Administration."),
    ("Drill down:", "tap any area to see how much was spent and where."),
    ("Dashboard:", "community stats + spending charts at a glance."),
    ("Sourced from:", "Ferrus ASAP Suite + Adagio / Sage 300 (when integrated)."),
], size=16, gap=1.4)
card(s, Inches(8.2), Inches(2.0), Inches(4.5), Inches(4.6), "Why it matters",
     "Transparent band management builds trust. Members can see budgets vs. actual "
     "spending across housing, education, health and more — the same financial "
     "system staff already use, surfaced openly in the app.", accent=SUCCESS)

# ---------- Slide 4: Architecture ----------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Architecture", "Azure end-to-end")
card(s, Inches(0.7), Inches(2.0), Inches(3.85), Inches(2.0), "Skintyee App",
     "React Native + Expo. iOS, Android, and web from one codebase. "
     "Material UI, role-based menus.", accent=PRIMARY)
card(s, Inches(4.74), Inches(2.0), Inches(3.85), Inches(2.0), "Microsoft Entra ID",
     "Identity & sign-in (Azure AD). App roles map to Public / Member / Admin. "
     "Replaces AWS Cognito.", accent=ACCENT)
card(s, Inches(8.78), Inches(2.0), Inches(3.85), Inches(2.0), "Azure Blob Storage",
     "Documents & media. Durable, backed-up. Replaces AWS S3.", accent=ACCENT)
card(s, Inches(0.7), Inches(4.2), Inches(3.85), Inches(2.0), "API Server",
     "App.SkinTyee.ca. Serves all app data; enforces role permissions.", accent=PRIMARY)
card(s, Inches(4.74), Inches(4.2), Inches(3.85), Inches(2.0), "Azure Cloud DB",
     "Managed database — automated backups + point-in-time restore.", accent=PRIMARY)
card(s, Inches(8.78), Inches(4.2), Inches(3.85), Inches(2.0), "Integrations",
     "Ferrus ASAP + Adagio / Sage 300 (finances); skintyee.ca WordPress "
     "(news & notification categories); auto-publish; push.", accent=SUCCESS)
txt(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.6),
    [("Azure-everywhere: matches the Nation's existing Microsoft footprint — Azure DNS, Outlook, Teams, Azure Storage.", 12.5, MUTED, False)])

# ---------- Slide 5: POC approach ----------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Approach", "Proof-of-concept first")
bullets(s, Inches(0.7), Inches(2.0), Inches(7.2), Inches(4.5), [
    ("POC goal:", "demonstrate the full experience and every role's menu — now."),
    ("Mock data:", "all screens render realistic content behind a clean data seam."),
    ("Role switcher:", "preview Public / Member / Admin instantly, no login needed."),
    ("Real services later:", "swap mocks for the live API, Entra ID, and Azure Blob"),
    ("", "with no rework of the screens — they only depend on the data contract."),
], size=17, gap=1.5)
card(s, Inches(8.2), Inches(2.0), Inches(4.5), Inches(4.4), "De-risked delivery",
     "Decisions, look-and-feel, and scope are validated up front with a working "
     "demo — before investing in backend infrastructure. The POC IS the first "
     "phase of the build, not throwaway work.", accent=SUCCESS)

# ---------- Slide 6: Roadmap ----------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Roadmap", "3-month engagement")
phases = [
    ("Phase 1 · POC", "May – mid-June",
     "RN app, theme & stack, all 7 features with role gating, mock data, demo-ready.", PRIMARY),
    ("Phase 2 · Backend & Identity", "mid-June – mid-July",
     "API Server + Azure Cloud DB live; Microsoft Entra ID sign-in; retire mocks.", ACCENT),
    ("Phase 3 · Storage, Integrations & Hardening", "mid-July – end Aug",
     "Azure Blob Storage, push, auto-publish, branding, testing, native builds, UAT.", SUCCESS),
]
y = Inches(2.1)
for t, when, body, c in phases:
    box(s, Inches(0.7), y, SW - Inches(1.4), Inches(1.45), fill=CARD,
        line=RGBColor(0x3A, 0x3A, 0x3A))
    box(s, Inches(0.7), y, Emu(int(Inches(0.1))), Inches(1.45), fill=c)
    txt(s, Inches(1.0), y + Inches(0.16), Inches(8.5), Inches(0.5), [(t, 18, TEXT, True)])
    txt(s, Inches(1.0), y + Inches(0.72), Inches(9.5), Inches(0.6), [(body, 13, MUTED, False)])
    txt(s, Inches(10.0), y + Inches(0.45), Inches(2.6), Inches(0.6),
        [(when, 14, c, True)], align=PP_ALIGN.RIGHT)
    y += Inches(1.65)

# ---------- Slide 7: Testing & distribution ----------
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "Quality", "Testing & distribution")
card(s, Inches(0.7), Inches(2.0), Inches(3.85), Inches(2.1), "Unit & integration",
     "Jest tests for the store, thunks, and mock API; React Native Testing "
     "Library for screens and role gating.", accent=PRIMARY)
card(s, Inches(4.74), Inches(2.0), Inches(3.85), Inches(2.1), "End-to-end",
     "Maestro / Detox for native flows; Cypress for web. Runs in Azure DevOps CI "
     "on every push.", accent=PRIMARY)
card(s, Inches(8.78), Inches(2.0), Inches(3.85), Inches(2.1), "UAT",
     "Band staff & council validate every feature per role against a scripted "
     "checklist before release.", accent=ACCENT)
card(s, Inches(0.7), Inches(4.3), Inches(5.9), Inches(2.0), "iOS — TestFlight",
     "EAS Build → TestFlight for internal & external testers on their own "
     "devices. Promote to App Store after UAT.", accent=ACCENT)
card(s, Inches(6.74), Inches(4.3), Inches(5.9), Inches(2.0), "Android — Google Play",
     "EAS Build → Google Play internal & closed testing tracks. Promote to "
     "production after UAT.", accent=SUCCESS)

# ---------- Slide 8: Closing ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0, 0, SW, Inches(0.22), fill=PRIMARY)
box(s, 0, Inches(0.22), Inches(5), Emu(int(Inches(0.08))), fill=ACCENT)
txt(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.2),
    [("A modern, self-owned platform", 40, TEXT, True)])
txt(s, Inches(0.92), Inches(3.8), Inches(11.5), Inches(0.8),
    [("for the Skintyee First Nation community.", 24, PRIMARY, True)])
txt(s, Inches(0.92), Inches(5.0), Inches(11.5), Inches(0.6),
    [("Working POC today · production-ready in three months.", 16, MUTED, False)])
txt(s, Inches(0.92), Inches(6.7), Inches(11.5), Inches(0.4),
    [("Blue Collar Dev  ·  lucas@bluecollardev.com", 13, ACCENT, True)])

out = "/Users/lucas/Workspaces/skintyee/docs/Skintyee-App-Proposal.pptx"
prs.save(out)
print("saved", out, "with", len(prs.slides._sldIdLst), "slides")
